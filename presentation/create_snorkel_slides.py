# -*- coding: utf-8 -*-
"""
車中泊×シュノーケル 10分プレゼン用 PowerPoint 生成スクリプト
構成C（大家族最適解）× A（王道ストーリー）完全版
画像プレースホルダー付き：挿入位置・画像ブリーフを明示し、違和感のないレイアウト
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from pptx.dml.color import RGBColor
except ImportError:
    RGBColor = None

# プレースホルダー用の薄い色（画像差し替え時に消す）
PLACEHOLDER_FILL_RGB = (240, 248, 255)  # 薄い水色
PLACEHOLDER_BORDER_RGB = (173, 216, 230)  # ライトブルー

def _apply_placeholder_style(shape, fill_rgb=PLACEHOLDER_FILL_RGB, border_rgb=PLACEHOLDER_BORDER_RGB):
    """図形にプレースホルダー用の塗りと枠を設定"""
    if RGBColor:
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
        except Exception:
            pass
        try:
            shape.line.color.rgb = RGBColor(*border_rgb)
        except Exception:
            pass

def add_title_slide(prs, title, subtitle="", image_brief=None):
    """タイトルスライド。image_brief があれば下部にメイン画像用のプレースホルダーを配置"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # タイトル・サブタイトル（上部）
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.0))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Yu Gothic"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.name = "Yu Gothic"
        p2.space_before = Pt(12)
    # 画像エリア（中央〜下部）：ここに画像を挿入すると違和感なし
    if image_brief:
        img_left, img_top = Inches(0.5), Inches(1.5)
        img_w, img_h = Inches(9), Inches(5.3)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_w, img_h)
        _apply_placeholder_style(rect)
        rect.text_frame.word_wrap = True
        rect.text_frame.paragraphs[0].text = image_brief
        rect.text_frame.paragraphs[0].font.size = Pt(11)
        rect.text_frame.paragraphs[0].font.name = "Yu Gothic"
        rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        rect.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return slide

def add_slide_with_image_placeholder(prs, title, bullets, image_brief, footer=""):
    """左に画像ゾーン・右にテキスト。画像挿入時に違和感のないレイアウト"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 左：画像プレースホルダー（ここに画像を挿入）
    img_left, img_top = Inches(0.5), Inches(1.0)
    img_w, img_h = Inches(4.4), Inches(5.2)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_w, img_h)
    _apply_placeholder_style(rect)
    rect.text_frame.word_wrap = True
    rect.text_frame.paragraphs[0].text = image_brief
    rect.text_frame.paragraphs[0].font.size = Pt(10)
    rect.text_frame.paragraphs[0].font.name = "Yu Gothic"
    rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    rect.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 右：タイトル＋箇条書き
    tx_left, tx_top = Inches(5.0), Inches(0.3)
    tx_w, tx_h = Inches(4.5), Inches(0.65)
    tbox = slide.shapes.add_textbox(tx_left, tx_top, tx_w, tx_h)
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
    bbox = slide.shapes.add_textbox(Inches(5.0), Inches(1.05), Inches(4.5), Inches(5.2))
    bf = bbox.text_frame
    bf.word_wrap = True
    for i, line in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = "• " + line if not line.startswith("•") else line
        para.font.size = Pt(16)
        para.font.name = "Yu Gothic"
        para.space_after = Pt(6)
    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.35), Inches(9), Inches(0.7))
        fp = fbox.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(14)
        fp.font.bold = True
        fp.font.name = "Yu Gothic"
        fp.alignment = PP_ALIGN.CENTER
    return slide

def add_slide_two_image_placeholders(prs, title, left_brief, right_brief, footer=""):
    """左右2枚の画像プレースホルダー（例：NG vs OK）。中央にタイトル、下に一言"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # タイトル
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
    # 左プレースホルダー
    w, h = Inches(4.2), Inches(4.0)
    r1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.0), w, h)
    _apply_placeholder_style(r1)
    r1.text_frame.word_wrap = True
    r1.text_frame.paragraphs[0].text = left_brief
    r1.text_frame.paragraphs[0].font.size = Pt(10)
    r1.text_frame.paragraphs[0].font.name = "Yu Gothic"
    r1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r1.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 右プレースホルダー
    r2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.0), w, h)
    _apply_placeholder_style(r2)
    r2.text_frame.word_wrap = True
    r2.text_frame.paragraphs[0].text = right_brief
    r2.text_frame.paragraphs[0].font.size = Pt(10)
    r2.text_frame.paragraphs[0].font.name = "Yu Gothic"
    r2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.5), Inches(5.15), Inches(9), Inches(0.7))
        fbox.text_frame.paragraphs[0].text = footer
        fbox.text_frame.paragraphs[0].font.size = Pt(14)
        fbox.text_frame.paragraphs[0].font.bold = True
        fbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
        fbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # タイトル
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Yu Gothic"
    # 本文（箇条書き）
    bbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(8.5), Inches(5))
    bf = bbox.text_frame
    bf.word_wrap = True
    for i, line in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = "• " + line if not line.startswith("•") else line
        para.font.size = Pt(18)
        para.font.name = "Yu Gothic"
        para.space_after = Pt(8)
    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
        ff = fbox.text_frame
        fp = ff.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(16)
        fp.font.bold = True
        fp.font.name = "Yu Gothic"
        fp.alignment = PP_ALIGN.CENTER
    return slide

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Yu Gothic"
    cols, col_width = len(headers), Inches(2.0)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.0), col_width * cols, Inches(2.2)).table
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        table.cell(0, c).text_frame.paragraphs[0].font.bold = True
        table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(14)
        table.cell(0, c).text_frame.paragraphs[0].font.name = "Yu Gothic"
    for r, row in enumerate(rows):
        for c, cell in enumerate(row):
            if c < cols:
                table.cell(r + 1, c).text = str(cell)
                table.cell(r + 1, c).text_frame.paragraphs[0].font.size = Pt(12)
                table.cell(r + 1, c).text_frame.paragraphs[0].font.name = "Yu Gothic"
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 【1】タイトル（下部にメイン画像用プレースホルダー）
    add_title_slide(prs,
        "家族8人でも最強に楽しめる\n車中泊×シュノーケル",
        "子ども6人・ハイエースで海へ",
        image_brief="【ここに画像を挿入】\n検索: 海 子ども 後ろ姿 家族\n縦横比: 16:9 推奨\n雰囲気: 明るい海、家族の一体感"
    )

    # 【2】導入
    add_content_slide(prs, "自己紹介＋導入", [
        "子ども6人の8人家族",
        "ハイエースで車中泊しながら海へ",
        "海での子どもたちのリアクション（感情描写）",
        "そこで出会ったのが「シュノーケル」でした。",
    ])

    # 【3】シュノーケルとは
    add_content_slide(prs, "シュノーケルとは", [
        "特別な資格不要",
        "浅瀬でOK",
        "子どもでもできる",
    ], "最も始めやすいマリンスポーツ")

    # 【4】マリンスポーツの全体像
    add_content_slide(prs, "マリンスポーツの全体像", [
        "【手軽ゾーン】シュノーケル / SUP / カヤック",
        "【本格潜水】スキューバダイビング / フリーダイビング",
        "【ライド系】サーフィン / ジェットスキー",
        "【体験レジャー】バナナボート / パラセーリング",
    ], "この中で、最も参入障壁が低いのがシュノーケルです。")

    # 【5】大家族視点での比較
    add_table_slide(prs, "大家族視点での比較",
        ["項目", "シュノーケル", "ダイビング", "サーフィン"],
        [
            ["初期費用", "◎ 安い", "△ 高い", "○"],
            ["技術習得", "◎ 不要", "△ 必須", "△"],
            ["年齢対応", "◎ 幅広い", "△ 制限あり", "△"],
            ["家族同時参加", "◎", "△", "×"],
        ]
    )

    # 【6】シュノーケルの立ち位置
    add_content_slide(prs, "シュノーケルの立ち位置", [
        "入口でありながら、装備次第で本格スポーツに",
        "参入難易度：最も低い / 装備コスト：最も安い",
        "年齢対応幅：最も広い / 家族適性：最強クラス",
    ], "「ただの海遊び」ではなく、最も裾野が広いマリンスポーツ")

    # 【7】名古屋から行けるおすすめ（導入）
    add_content_slide(prs, "名古屋から行けるおすすめ4選", [
        "① 水晶浜（福井）— ファミリー入門",
        "② 水島（福井）— 北陸のハワイ",
        "③ ヒリゾ浜（静岡）— 本州トップ級の透明度",
        "④ 串本（和歌山）— 世界最北のサンゴ礁 ※実体験厚め",
    ])

    # 【8】水晶浜（左に画像ゾーン）
    add_slide_with_image_placeholder(prs, "① 水晶浜（福井）", [
        "遠浅で子連れ最強",
        "透明度が高い",
        "アクセス良好",
    ], "【画像】\n検索: 水晶浜 福井 海\n比: 4:5\n雰囲気: 遠浅・透明", "→ ファミリー入門に最適")

    # 【9】水島（左に画像ゾーン）
    add_slide_with_image_placeholder(prs, "② 水島（福井）", [
        "船で渡る特別感",
        "「北陸のハワイ」",
        "子どもテンションMAX",
    ], "【画像】\n検索: 水島 福井 船 海岸\n比: 4:5\n雰囲気: 非日常・南国", "→ 非日常体験ができる")

    # 【10】ヒリゾ浜（左に画像ゾーン）
    add_slide_with_image_placeholder(prs, "③ ヒリゾ浜（静岡）", [
        "本州トップ級の透明度",
        "魚の量が段違い",
        "やや上級者向け",
    ], "【画像】\n検索: ヒリゾ浜 静岡 透明度\n比: 4:5\n雰囲気: 青く澄んだ海", "→ 感動体験枠")

    # 【11】串本（左に画像ゾーン・あなたの写真推奨）
    add_slide_with_image_placeholder(prs, "④ 串本（和歌山）※最重要", [
        "世界最北のサンゴ礁",
        "生き物が豊富",
        "家族旅行との相性◎",
        "実際の子どもたちの反応が非常に大きかった",
    ], "【画像・あなたの写真推奨】\n検索: 串本 サンゴ 和歌山\n比: 4:5\n雰囲気: 子どもの笑顔・感動", "→ あなたの実体験エピソードを厚く")

    # 【12】シュノーケルの楽しみ方
    add_content_slide(prs, "シュノーケルの楽しみ方", [
        "魚を探す＝海の宝探し感覚",
        "親子で同じ景色を共有",
        "写真・動画で思い出倍増",
        "朝イチが最強（車中泊と相性◎）",
    ])

    # 【13】安全装備の優先順位
    add_content_slide(prs, "安全に楽しむための装備優先順位", [
        "【最優先】ライフジャケット / マリンシューズ / マスク＆シュノーケル",
        "【快適性アップ】フィン / ドライシュノーケル / 度付きゴーグル",
        "【余裕があれば】マリン手袋 / ラッシュガード",
    ], "道具を整えると「遊び」から「スポーツ」に変わります。")

    # 【14】マリンシューズ＋スポーツ装備
    add_content_slide(prs, "最重要：マリンシューズ＆本格装備", [
        "マリンシューズ：足のケガ防止・滑り止め・フィン擦れ防止（一番地味で一番重要）",
        "フィン：股関節から大きくゆっくり。自転車こぎはNG",
        "ドライシュノーケル：水が入りにくく子ども向き",
        "度付きゴーグル：見え方が変わると感動の量が変わる",
        "マリン専用手袋：岩場・滑り止め・ケガ防止",
    ])

    # 【15】フィンの正しいキック図解（左NG・右OKの画像2枚）
    add_slide_two_image_placeholders(prs, "フィンの正しい使い方",
        "【左に画像】\n❌ NGキック\n検索: 自転車こぎ キック 水泳\n比: 1:1\n膝だけ・バシャバシャ",
        "【右に画像】\n✅ 正しいキック\n検索: フィン キック 股関節\n比: 1:1\nゆったり大きく",
        "フィンは「速く」ではなく「ゆっくり大きく」が正解です。"
    )

    # 【16】車中泊×シュノーケル
    add_content_slide(prs, "車中泊×シュノーケル最強説", [
        "海の近くに前泊できる → 朝イチで一番きれいな海に入れる",
        "宿泊費が浮く",
        "「海遊びの満足度は、朝の一番風呂で決まります。」",
    ])

    # 【17】ハイエース車内装備（左に画像ゾーン）
    add_slide_with_image_placeholder(prs, "ハイエース車内装備", [
        "サブバッテリー（電源の心配なし）",
        "電子レンジ（コンビニ飯が温かい・雨の日の神）",
        "テレビ＋Amazon Fire TV Stick（夜は動く映画館）",
        "キッチンセット（朝ごはんを現地で作れる）",
        "炊飯器（海上がりの炊きたては反則級）",
    ], "【画像・あなたの写真推奨】\n検索: ハイエース 車中泊 内装\n比: 4:5\n雰囲気: テレビ・キッチン・くつろぎ", "もはや小さな動く家です。")

    # 【18】他レジャーとのコスパ比較
    add_table_slide(prs, "他レジャーとのコスパ比較",
        ["項目", "シュノーケル", "遊園地", "キャンプ"],
        [
            ["初期費用", "◎ 安い", "なし", "△ 高い"],
            ["維持費", "ほぼゼロ", "毎回高い", "かかる"],
            ["子ども満足度", "◎", "◎", "○"],
            ["自然体験", "◎", "△", "◎"],
        ]
    )
    slide = prs.slides[-1]
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.5))
    tbox.text_frame.paragraphs[0].text = "→ 大家族ほど最強の遊び"
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.size = Pt(16)
    tbox.text_frame.paragraphs[0].font.name = "Yu Gothic"

    # 【19】まとめ（感動締め）
    add_content_slide(prs, "まとめ", [
        "豪華なホテルに泊まらなくても、高価なアクティビティがなくても、",
        "家族みんなで同じ海をのぞいた時間は、想像以上に心に残ります。",
        "子どもが大きくなっても、一緒に同じ海をのぞいた記憶は残り続けます。",
    ], "次の休日、ちょっと海をのぞいてみませんか？")

    import os
    out_path = os.path.join(os.path.dirname(__file__), "車中泊×シュノーケル_画像プレースホルダー版.pptx")
    prs.save(out_path)
    print(f"保存しました: {out_path}")

if __name__ == "__main__":
    main()
