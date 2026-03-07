# -*- coding: utf-8 -*-
"""
南の海テンプレ風 × 提供いただいた4枚の画像で車中泊×シュノーケルPPTを再現
- タイトル：全面背景画像＋半透明円形にタイトル
- コンテンツ：左に画像 or 青バー、右にテキスト
- 表：ネイビーヘッダー・南の海カラー
- ラスト：THANK YOU 全面画像
"""
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from pptx.dml.color import RGBColor
except ImportError:
    RGBColor = None

# 南の海パレット（参照テンプレに合わせる）
NAVY = RGBColor(0x1B, 0x4F, 0x72) if RGBColor else None
LIGHT_BLUE = RGBColor(0x87, 0xCE, 0xEB) if RGBColor else None  # 薄い青オーバーレイ
YELLOW_ACCENT = RGBColor(0xFF, 0xD7, 0x00) if RGBColor else None
WHITE = RGBColor(0xFF, 0xFF, 0xFF) if RGBColor else None
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50) if RGBColor else None


def find_beach_images(images_dir, uuids=("c72406ed", "056ad239", "e3f6f5b4", "e5f74cef")):
    """imagesフォルダから指定UUIDを含むPNGを順に返す（最大4枚）"""
    if not os.path.isdir(images_dir):
        return []
    found = []
    for uid in uuids:
        for path in glob.glob(os.path.join(images_dir, "*" + uid + "*.png")):
            found.append(path)
            break
    return found[:4]


def add_background_picture(slide, prs, image_path, width=None, height=None):
    """スライド全面に背景画像を追加（最背面になるよう先に追加）"""
    if not image_path or not os.path.isfile(image_path):
        return
    w = width if width is not None else Inches(10)
    h = height if height is not None else Inches(7.5)
    slide.shapes.add_picture(image_path, Inches(0), Inches(0), w, h)


def add_overlay_rect(slide, left, top, width, height, rgb, transparency=0.4):
    """半透明のオーバーレイ四角（テキスト用）"""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if RGBColor:
        rect.fill.solid()
        rect.fill.fore_color.rgb = rgb
        try:
            rect.fill.fore_color.brightness = -transparency
        except Exception:
            pass
        try:
            rect.fill.transparency = transparency
        except Exception:
            pass
        rect.line.fill.background()
    return rect


def add_title_slide_south_sea(prs, image_path, title_line1, title_line2, subtitle):
    """タイトルスライド：全面背景＋右に半透明円形でタイトル（南の海テンプレ風）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_picture(slide, prs, image_path)
    # 右側に大きな半透明円（角丸四角で代用）
    cx, cy = Inches(5.5), Inches(3.5)
    r = Inches(2.8)
    left = cx - r
    top = cy - r
    ov = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, r * 2, r * 2)
    if RGBColor:
        ov.fill.solid()
        ov.fill.fore_color.rgb = LIGHT_BLUE
        try:
            ov.fill.transparency = 0.35
        except Exception:
            pass
        ov.line.fill.background()
    # タイトル（円内）
    tb = slide.shapes.add_textbox(Inches(4.0), Inches(2.2), Inches(4.5), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title_line1
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.name = "Yu Gothic"
    if RGBColor:
        p1.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = title_line2
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.name = "Yu Gothic"
    if RGBColor:
        p2.font.color.rgb = YELLOW_ACCENT
    p2.space_before = Pt(6)
    p3 = tf.add_paragraph()
    p3.text = subtitle
    p3.font.size = Pt(14)
    p3.font.name = "Yu Gothic"
    if RGBColor:
        p3.font.color.rgb = WHITE
    p3.space_before = Pt(10)
    return slide


def add_content_with_left_image(prs, image_path, title, bullets, footer=""):
    """左に画像・右に青バー風＋テキスト（セクションスライド風）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if image_path and os.path.isfile(image_path):
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(4.5), prs.slide_height)
    # 右側コンテンツエリアに薄い青の縦バー（CONTENTS風）
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(0), Inches(0.4), prs.slide_height)
    if RGBColor:
        bar.fill.solid()
        bar.fill.fore_color.rgb = LIGHT_BLUE
        try:
            bar.fill.transparency = 0.5
        except Exception:
            pass
        bar.line.fill.background()
    # タイトル
    tbox = slide.shapes.add_textbox(Inches(5.0), Inches(0.35), Inches(4.5), Inches(0.7))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
    if RGBColor:
        tbox.text_frame.paragraphs[0].font.color.rgb = NAVY
    # 箇条書き
    bbox = slide.shapes.add_textbox(Inches(5.0), Inches(1.15), Inches(4.5), Inches(5.2))
    bf = bbox.text_frame
    bf.word_wrap = True
    for i, line in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = "• " + line if not line.startswith("•") else line
        p.font.size = Pt(15)
        p.font.name = "Yu Gothic"
        if RGBColor:
            p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)
    if footer:
        fbox = slide.shapes.add_textbox(Inches(5.0), Inches(6.4), Inches(4.5), Inches(0.6))
        fbox.text_frame.paragraphs[0].text = footer
        fbox.text_frame.paragraphs[0].font.size = Pt(13)
        fbox.text_frame.paragraphs[0].font.bold = True
        fbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
        if RGBColor:
            fbox.text_frame.paragraphs[0].font.color.rgb = NAVY
    return slide


def add_content_slide_with_header(prs, image_path, title, bullets, footer=""):
    """上部にグラデーション風ヘッダーバー＋背景画像（オプション）＋本文"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if image_path and os.path.isfile(image_path):
        add_background_picture(slide, prs, image_path)
        # 薄いオーバーレイで文字を見やすく
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        if RGBColor:
            ov.fill.solid()
            ov.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                ov.fill.transparency = 0.6
            except Exception:
                pass
            ov.line.fill.background()
    # ヘッダーバー（薄い青）
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    if RGBColor:
        bar.fill.solid()
        bar.fill.fore_color.rgb = LIGHT_BLUE
        try:
            bar.fill.transparency = 0.3
        except Exception:
            pass
        bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(26)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
    if RGBColor:
        tbox.text_frame.paragraphs[0].font.color.rgb = NAVY
    bbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9), Inches(5.5))
    bf = bbox.text_frame
    bf.word_wrap = True
    for i, line in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = "• " + line if not line.startswith("•") else line
        p.font.size = Pt(16)
        p.font.name = "Yu Gothic"
        if RGBColor:
            p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(8)
    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.6))
        fbox.text_frame.paragraphs[0].text = footer
        fbox.text_frame.paragraphs[0].font.size = Pt(14)
        fbox.text_frame.paragraphs[0].font.bold = True
        fbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        fbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
        if RGBColor:
            fbox.text_frame.paragraphs[0].font.color.rgb = NAVY
    return slide


def add_table_slide_south_sea(prs, title, headers, rows):
    """表スライド（南の海：ネイビーヘッダー・白/サンド body）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.75))
    if RGBColor:
        bar.fill.solid()
        bar.fill.fore_color.rgb = LIGHT_BLUE
        try:
            bar.fill.transparency = 0.3
        except Exception:
            pass
        bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    tbox.text_frame.paragraphs[0].text = title
    tbox.text_frame.paragraphs[0].font.size = Pt(24)
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
    if RGBColor:
        tbox.text_frame.paragraphs[0].font.color.rgb = NAVY
    cols = len(headers)
    col_w = Inches(2.0)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.0), col_w * cols, Inches(2.0)).table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.name = "Yu Gothic"
        if RGBColor:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if c >= cols:
                continue
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.name = "Yu Gothic"
            if RGBColor:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF4, 0xE4, 0xD0) if r % 2 == 0 else WHITE
                cell.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
    return slide


def add_thank_you_slide(prs, image_path):
    """ラスト：全面背景＋中央に THANK YOU"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background_picture(slide, prs, image_path)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.0))
    tf = tb.text_frame
    tf.paragraphs[0].text = "THANK YOU"
    tf.paragraphs[0].font.size = Pt(56)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Yu Gothic"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if RGBColor:
        tf.paragraphs[0].font.color.rgb = WHITE
    return slide


def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    images_dir = os.path.join(base_dir, "images")
    images = find_beach_images(images_dir)
    if len(images) < 4:
        # フォルダ直下のpngを先頭4枚使うフォールバック
        all_png = glob.glob(os.path.join(images_dir, "*.png"))
        images = (all_png + [None] * 4)[:4]
    img1 = images[0] if len(images) > 0 else None
    img2 = images[1] if len(images) > 1 else None
    img3 = images[2] if len(images) > 2 else None
    img4 = images[3] if len(images) > 3 else None

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 【1】タイトル（全面画像＋円形タイトル）
    add_title_slide_south_sea(prs, img1,
        "家族8人でも最強に楽しめる",
        "車中泊×シュノーケル",
        "子ども6人・ハイエースで海へ"
    )

    # 【2】導入
    add_content_slide_with_header(prs, None, "自己紹介＋導入", [
        "子ども6人の8人家族",
        "ハイエースで車中泊しながら海へ",
        "海での子どもたちのリアクション",
        "そこで出会ったのが「シュノーケル」でした。",
    ])

    # 【3】シュノーケルとは
    add_content_slide_with_header(prs, None, "シュノーケルとは", [
        "特別な資格不要",
        "浅瀬でOK",
        "子どもでもできる",
    ], "最も始めやすいマリンスポーツ")

    # 【4】マリンスポーツの全体像
    add_content_slide_with_header(prs, None, "マリンスポーツの全体像", [
        "【手軽ゾーン】シュノーケル / SUP / カヤック",
        "【本格潜水】スキューバダイビング / フリーダイビング",
        "【ライド系】サーフィン / ジェットスキー",
        "【体験レジャー】バナナボート / パラセーリング",
    ], "この中で、最も参入障壁が低いのがシュノーケルです。")

    # 【5】大家族視点での比較（表）
    add_table_slide_south_sea(prs, "大家族視点での比較",
        ["項目", "シュノーケル", "ダイビング", "サーフィン"],
        [
            ["初期費用", "◎ 安い", "△ 高い", "○"],
            ["技術習得", "◎ 不要", "△ 必須", "△"],
            ["年齢対応", "◎ 幅広い", "△ 制限あり", "△"],
            ["家族同時参加", "◎", "△", "×"],
        ]
    )

    # 【6】シュノーケルの立ち位置
    add_content_slide_with_header(prs, None, "シュノーケルの立ち位置", [
        "入口でありながら、装備次第で本格スポーツに",
        "参入難易度：最も低い / 装備コスト：最も安い",
        "年齢対応幅：最も広い / 家族適性：最強クラス",
    ], "「ただの海遊び」ではなく、最も裾野が広いマリンスポーツ")

    # 【7】おすすめ4選
    add_content_slide_with_header(prs, None, "名古屋から行けるおすすめ4選", [
        "① 水晶浜（福井）— ファミリー入門",
        "② 水島（福井）— 北陸のハワイ",
        "③ ヒリゾ浜（静岡）— 本州トップ級の透明度",
        "④ 串本（和歌山）— 世界最北のサンゴ礁",
    ])

    # 【8】【9】【10】【11】スポット紹介（左に画像）
    add_content_with_left_image(prs, img1, "① 水晶浜（福井）", [
        "遠浅で子連れ最強", "透明度が高い", "アクセス良好",
    ], "→ ファミリー入門に最適")
    add_content_with_left_image(prs, img2, "② 水島（福井）", [
        "船で渡る特別感", "「北陸のハワイ」", "子どもテンションMAX",
    ], "→ 非日常体験ができる")
    add_content_with_left_image(prs, img3, "③ ヒリゾ浜（静岡）", [
        "本州トップ級の透明度", "魚の量が段違い", "やや上級者向け",
    ], "→ 感動体験枠")
    add_content_with_left_image(prs, img4, "④ 串本（和歌山）※最重要", [
        "世界最北のサンゴ礁", "生き物が豊富", "家族旅行との相性◎",
        "実際の子どもたちの反応が非常に大きかった",
    ], "→ あなたの実体験エピソードを厚く")

    # 【12】楽しみ方
    add_content_slide_with_header(prs, img2, "シュノーケルの楽しみ方", [
        "魚を探す＝海の宝探し感覚",
        "親子で同じ景色を共有",
        "写真・動画で思い出倍増",
        "朝イチが最強（車中泊と相性◎）",
    ])

    # 【13】安全装備
    add_content_slide_with_header(prs, None, "安全に楽しむための装備優先順位", [
        "【最優先】ライフジャケット / マリンシューズ / マスク＆シュノーケル",
        "【快適性アップ】フィン / ドライシュノーケル / 度付きゴーグル",
        "【余裕があれば】マリン手袋 / ラッシュガード",
    ], "道具を整えると「遊び」から「スポーツ」に変わります。")

    # 【14】マリンシューズ＋本格装備
    add_content_slide_with_header(prs, None, "最重要：マリンシューズ＆本格装備", [
        "マリンシューズ：足のケガ防止・滑り止め・フィン擦れ防止",
        "フィン：股関節から大きくゆっくり。自転車こぎはNG",
        "ドライシュノーケル：水が入りにくく子ども向き",
        "度付きゴーグル：見え方が変わると感動の量が変わる",
        "マリン専用手袋：岩場・滑り止め・ケガ防止",
    ])

    # 【15】フィンの正しいキック
    add_content_slide_with_header(prs, None, "フィンの正しい使い方", [
        "❌ NG：自転車こぎキック（膝だけ・水面を叩く・すぐ疲れる）",
        "✅ OK：股関節から大きくゆっくり（少ない力で長く泳げる）",
    ], "フィンは「速く」ではなく「ゆっくり大きく」が正解です。")

    # 【16】車中泊
    add_content_slide_with_header(prs, img3, "車中泊×シュノーケル最強説", [
        "海の近くに前泊できる → 朝イチで一番きれいな海に入れる",
        "宿泊費が浮く",
        "「海遊びの満足度は、朝の一番風呂で決まります。」",
    ])

    # 【17】ハイエース車内装備
    add_content_slide_with_header(prs, None, "ハイエース車内装備", [
        "サブバッテリー（電源の心配なし）",
        "電子レンジ（コンビニ飯が温かい・雨の日の神）",
        "テレビ＋Amazon Fire TV Stick（夜は動く映画館）",
        "キッチンセット（朝ごはんを現地で作れる）",
        "炊飯器（海上がりの炊きたては反則級）",
    ], "もはや小さな動く家です。")

    # 【18】コスパ比較表
    add_table_slide_south_sea(prs, "他レジャーとのコスパ比較",
        ["項目", "シュノーケル", "遊園地", "キャンプ"],
        [
            ["初期費用", "◎ 安い", "なし", "△ 高い"],
            ["維持費", "ほぼゼロ", "毎回高い", "かかる"],
            ["子ども満足度", "◎", "◎", "○"],
            ["自然体験", "◎", "△", "◎"],
        ]
    )
    slide = prs.slides[-1]
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(3.15), Inches(9), Inches(0.45))
    tbox.text_frame.paragraphs[0].text = "→ 大家族ほど最強の遊び"
    tbox.text_frame.paragraphs[0].font.bold = True
    tbox.text_frame.paragraphs[0].font.size = Pt(14)
    tbox.text_frame.paragraphs[0].font.name = "Yu Gothic"
    if RGBColor:
        tbox.text_frame.paragraphs[0].font.color.rgb = NAVY

    # 【19】まとめ（全面画像＋メッセージ）
    add_content_slide_with_header(prs, img4, "まとめ", [
        "豪華なホテルに泊まらなくても、高価なアクティビティがなくても、",
        "家族みんなで同じ海をのぞいた時間は、想像以上に心に残ります。",
        "子どもが大きくなっても、一緒に同じ海をのぞいた記憶は残り続けます。",
    ], "次の休日、ちょっと海をのぞいてみませんか？")

    # 【20】THANK YOU（全面画像）
    add_thank_you_slide(prs, img1)

    out_path = os.path.join(base_dir, "車中泊×シュノーケル_南の海スタイル.pptx")
    prs.save(out_path)
    print(f"保存しました: {out_path}")
    print(f"使用画像: {len(images)}枚")


if __name__ == "__main__":
    main()
