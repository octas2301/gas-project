# -*- coding: utf-8 -*-
"""
既存の .pptx に「南の海」テーマの装飾を適用するスクリプト
- トンマナ統一（カラー・フォントの一貫性）
- 表のスタイル（デフォルト表をやめてヘッダー・枠線を整える）
- タイトル下アクセント線・スライド内バランス
※画像差し込み済みのファイルを読み込み、装飾だけを上乗せして別名保存します。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE

try:
    from pptx.dml.color import RGBColor
except ImportError:
    RGBColor = None

if not RGBColor:
    print("RGBColor が使えません。装飾の色付けはスキップします。")
    exit(1)

# ========== 南の海テーマ カラーパレット（トンマナ統一） ==========
NAVY = RGBColor(0x1B, 0x4F, 0x72)        # 深い海・タイトル・表ヘッダー
TURQUOISE = RGBColor(0x17, 0xA2, 0xB8)  # ターコイズ・アクセント
CORAL = RGBColor(0xE0, 0x7A, 0x5F)       # コーラル・アクセント線
SAND = RGBColor(0xF4, 0xE4, 0xD0)        # サンド・表の偶数行（薄い）
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_SEA = RGBColor(0xE8, 0xF6, 0xF8)  # 薄い海色・表背景
BORDER_LIGHT = RGBColor(0xB0, 0xD4, 0xE0)  # 表の枠線（薄いターコイズ）
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)  # 本文用の濃いグレー


def style_table(table):
    """表を南の海テーマでスタイル（ヘッダー・枠線・バランス）"""
    rows, cols = len(table.rows), len(table.columns)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            try:
                cell.fill.solid()
                if r == 0:
                    cell.fill.fore_color.rgb = NAVY
                    for para in cell.text_frame.paragraphs:
                        para.font.bold = True
                        para.font.size = Pt(13)
                        para.font.color.rgb = WHITE
                        para.font.name = "Yu Gothic"
                else:
                    cell.fill.fore_color.rgb = SAND if r % 2 == 0 else WHITE
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(11)
                        para.font.color.rgb = TEXT_DARK
                        para.font.name = "Yu Gothic"
            except Exception:
                pass
            try:
                cell.text_frame.paragraphs[0].alignment = -1  # 左揃えのまま
            except Exception:
                pass
    # 枠線を薄い色で統一（セル境界を見やすく）
    try:
        _set_table_borders_oxml(table, "B0D4E0")
    except Exception:
        pass


def _set_table_borders_oxml(table, color_hex="B0D4E0", width_emu=12700):
    """表の全セルに薄い枠線を付ける（既存の枠があれば追加しない）"""
    from pptx.oxml import parse_xml
    for row in table.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            existing = {e.tag.split("}")[-1] for e in tcPr if hasattr(e, "tag")}
            for tag in ["lnL", "lnR", "lnT", "lnB"]:
                if tag in existing:
                    continue
                ln = parse_xml(
                    f'<a:{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="{width_emu}">'
                    f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill></a:{tag}>'
                )
                tcPr.append(ln)
    return


def add_title_accent_line(slide, top_inches=0.92, width_inches=9, height_pt=3):
    """タイトル下にアクセント線（南の海テーマのコーラル）を追加"""
    left = Inches(0.5)
    top = Inches(top_inches)
    width = Inches(width_inches)
    height = Inches(height_pt / 72.0)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = CORAL
    try:
        line.line.fill.background()
    except Exception:
        try:
            line.line.width = Pt(0)
        except Exception:
            pass
    return line


def apply_text_theme(shape):
    """テキストに南の海テーマの色を適用（タイトル＝NAVY、本文＝TEXT_DARK）"""
    if not shape.has_text_frame:
        return
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            if not para.text.strip():
                continue
            font = para.font
            font.name = "Yu Gothic"
            try:
                size = font.size
                if size is not None and size.pt >= 24:
                    font.color.rgb = NAVY
                elif size is not None and size.pt >= 14 and size.pt < 24:
                    font.color.rgb = NAVY
                else:
                    font.color.rgb = TEXT_DARK
            except Exception:
                font.color.rgb = TEXT_DARK
    except Exception:
        pass


def apply_slide_decoration(slide, slide_idx):
    """1スライド分の装飾：アクセント線・テキスト色・表スタイル・バランス"""
    # タイトル下アクセント線（南の海のコーラル・全スライドでトンマナ統一）
    try:
        add_title_accent_line(slide, top_inches=0.92, width_inches=9, height_pt=3)
    except Exception:
        pass

    # 既存のシェイプにテーマ色・表スタイルを適用
    for shape in slide.shapes:
        try:
            if shape.has_table:
                style_table(shape.table)
            elif shape.has_text_frame:
                apply_text_theme(shape)
        except Exception:
            continue


def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    input_path = os.path.join(base_dir, "車中泊×シュノーケル_画像プレースホルダー版.pptx")
    output_path = os.path.join(base_dir, "車中泊×シュノーケル_装飾版.pptx")

    if not os.path.exists(input_path):
        print(f"入力ファイルが見つかりません: {input_path}")
        return

    prs = Presentation(input_path)

    for idx, slide in enumerate(prs.slides):
        try:
            apply_slide_decoration(slide, idx)
        except Exception as e:
            print(f"スライド {idx + 1} でエラー: {e}")

    prs.save(output_path)
    print(f"装飾を適用して保存しました: {output_path}")


if __name__ == "__main__":
    main()
